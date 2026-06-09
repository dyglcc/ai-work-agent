"""文件生成服务：PPT、Word、图表、图片"""
from __future__ import annotations

import base64
import io
import logging
import textwrap
from typing import Any

import httpx
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import matplotlib
matplotlib.use('Agg')  # 非 GUI 后端
import matplotlib.pyplot as plt

from app.config import settings

logger = logging.getLogger(__name__)


def generate_pptx(title: str, slides_data: list[dict[str, Any]]) -> bytes:
    """
    生成 PowerPoint 文件

    Args:
        title: PPT 标题
        slides_data: 幻灯片数据列表，每个元素包含 {"title": str, "content": list[str]}

    Returns:
        .pptx 文件的二进制数据
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    # 内容页
    for slide_data in slides_data:
        bullet_slide_layout = prs.slide_layouts[1]  # 标题 + 内容布局
        slide = prs.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_data.get("title", "")

        tf = body_shape.text_frame
        content_items = slide_data.get("content", [])

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)

    # 保存到内存
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_docx(title: str, sections: list[dict[str, Any]]) -> bytes:
    """
    生成 Word 文档

    Args:
        title: 文档标题
        sections: 章节数据列表，每个元素包含 {"heading": str, "paragraphs": list[str]}

    Returns:
        .docx 文件的二进制数据
    """
    doc = Document()

    # 添加标题
    heading = doc.add_heading(title, level=0)
    heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    # 添加章节
    for section in sections:
        section_heading = section.get("heading", "")
        if section_heading:
            doc.add_heading(section_heading, level=1)

        paragraphs = section.get("paragraphs", [])
        for para_text in paragraphs:
            p = doc.add_paragraph(para_text)
            p.style = 'Normal'

    # 保存到内存
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generate_chart(chart_type: str, data: dict[str, Any], title: str) -> bytes:
    """
    生成数据图表（PNG）

    Args:
        chart_type: 图表类型 (bar, line, pie)
        data: 图表数据 {"labels": [...], "values": [...]}
        title: 图表标题

    Returns:
        PNG 图片的二进制数据
    """
    plt.figure(figsize=(10, 6))
    plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
    plt.rcParams['axes.unicode_minus'] = False

    labels = data.get("labels", [])
    values = data.get("values", [])

    if chart_type == "bar":
        plt.bar(labels, values)
        plt.xlabel("类别")
        plt.ylabel("数值")
    elif chart_type == "line":
        plt.plot(labels, values, marker='o')
        plt.xlabel("类别")
        plt.ylabel("数值")
        plt.grid(True, alpha=0.3)
    elif chart_type == "pie":
        plt.pie(values, labels=labels, autopct='%1.1f%%')
    else:
        # 默认柱状图
        plt.bar(labels, values)

    plt.title(title, fontsize=14, fontweight='bold')
    plt.tight_layout()

    # 保存到内存
    buffer = io.BytesIO()
    plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
    plt.close()
    buffer.seek(0)
    return buffer.getvalue()


def _hex_to_rgb(value: str) -> tuple[float, float, float]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) / 255 for i in (0, 2, 4))


def _palette_colors(palette: str) -> dict[str, str]:
    palettes = {
        "blue": {"top": "#0f172a", "bottom": "#2563eb", "accent": "#38bdf8", "text": "#ffffff", "muted": "#dbeafe"},
        "gold": {"top": "#4a2d00", "bottom": "#f59e0b", "accent": "#fde68a", "text": "#fff7ed", "muted": "#ffedd5"},
        "red": {"top": "#450a0a", "bottom": "#dc2626", "accent": "#fbbf24", "text": "#fff7ed", "muted": "#fee2e2"},
        "green": {"top": "#052e16", "bottom": "#16a34a", "accent": "#bbf7d0", "text": "#f0fdf4", "muted": "#dcfce7"},
        "purple": {"top": "#2e1065", "bottom": "#7c3aed", "accent": "#f0abfc", "text": "#faf5ff", "muted": "#ede9fe"},
        "dark": {"top": "#020617", "bottom": "#334155", "accent": "#94a3b8", "text": "#f8fafc", "muted": "#e2e8f0"},
        "warm": {"top": "#431407", "bottom": "#ea580c", "accent": "#fed7aa", "text": "#fff7ed", "muted": "#ffedd5"},
    }
    return palettes.get((palette or "").lower(), palettes["blue"])


def _draw_decorative_shapes(ax: Any, style: str, colors: dict[str, str]) -> None:
    """根据风格绘制装饰性图形元素."""
    if style in ("科技", "tech", "modern", "商务", "business"):
        # 科技风格：几何线条和网格
        for i in range(5):
            y = 0.15 + i * 0.18
            ax.plot([0.05, 0.95], [y, y], transform=ax.transAxes,
                    color=colors["accent"], lw=0.5, alpha=0.15)
        for angle in [0, 45, 90, 135]:
            rad = np.radians(angle)
            cx, cy = 0.85, 0.85
            r = 0.12
            ax.add_patch(plt.Polygon(
                [(cx + r * np.cos(rad + i * np.pi / 2), cy + r * np.sin(rad + i * np.pi / 2))
                 for i in range(4)],
                transform=ax.transAxes, fill=False, ec=colors["accent"], lw=1.2, alpha=0.3
            ))
    elif style in ("国风", "chinese", "传统", "古典"):
        # 国风：云纹、圆形装饰
        for x, y, r in [(0.12, 0.82, 0.06), (0.88, 0.78, 0.05), (0.15, 0.25, 0.04)]:
            ax.add_patch(plt.Circle((x, y), r, transform=ax.transAxes,
                                    fill=False, ec=colors["accent"], lw=1.5, alpha=0.4))
            ax.add_patch(plt.Circle((x, y), r * 0.6, transform=ax.transAxes,
                                    fill=True, fc=colors["accent"], alpha=0.08, lw=0))
    elif style in ("温馨", "warm", "家庭", "family"):
        # 温馨：心形、圆点装饰
        for x, y, s in [(0.1, 0.85, 80), (0.9, 0.82, 60), (0.85, 0.2, 50)]:
            ax.scatter([x], [y], s=s, c=colors["accent"], alpha=0.2, transform=ax.transAxes)
    elif style in ("节日", "festival", "喜庆", "celebration"):
        # 节日：星星、彩带
        for x, y in [(0.1, 0.9), (0.9, 0.88), (0.88, 0.15), (0.12, 0.18)]:
            ax.add_patch(plt.RegularPolygon((x, y), 5, radius=0.03, transform=ax.transAxes,
                                            fill=True, fc=colors["accent"], alpha=0.25, lw=0))
    else:
        # 默认：光斑装饰
        for x, y, r, alpha in [(0.15, 0.88, 0.20, 0.18), (0.92, 0.72, 0.28, 0.16), (0.08, 0.22, 0.24, 0.13)]:
            ax.add_patch(plt.Circle((x, y), r, transform=ax.transAxes,
                                    color=colors["accent"], alpha=alpha, lw=0))


def _get_layout_template(style: str, colors: dict[str, str]) -> dict:
    """根据风格返回布局参数."""
    layouts = {
        "default": {
            "tag_y": 0.86, "line_y": 0.83,
            "title_y": 0.72, "subtitle_y": 0.60,
            "elements_start_y": 0.50, "element_spacing": 0.055,
            "body_y": 0.22, "brand_y": 0.10,
            "title_size": 43, "subtitle_size": 21, "body_size": 17,
        },
        "商务": {
            "tag_y": 0.88, "line_y": 0.85,
            "title_y": 0.75, "subtitle_y": 0.63,
            "elements_start_y": 0.52, "element_spacing": 0.05,
            "body_y": 0.25, "brand_y": 0.12,
            "title_size": 40, "subtitle_size": 19, "body_size": 16,
        },
        "国风": {
            "tag_y": 0.84, "line_y": 0.81,
            "title_y": 0.68, "subtitle_y": 0.56,
            "elements_start_y": 0.45, "element_spacing": 0.05,
            "body_y": 0.20, "brand_y": 0.08,
            "title_size": 46, "subtitle_size": 22, "body_size": 18,
        },
    }
    for key in layouts:
        if key in style:
            return layouts[key]
    return layouts["default"]


def generate_prompt_image(
    prompt: str,
    title: str = "AI 作图",
    subtitle: str = "",
    body: str = "",
    style: str = "",
    palette: str = "blue",
    elements: list[str] | None = None,
) -> bytes:
    """生成本地海报 PNG 兜底图，用于未配置外部生图 API 的场景.
    
    升级版本：更专业的设计、更好的排版、更多装饰元素、更高分辨率
    """
    colors = _palette_colors(palette)
    elements = [str(e) for e in (elements or []) if str(e).strip()][:6]
    subtitle = subtitle or style or "智能设计生成"
    body = body or "根据你的需求生成主题相关的海报预览，可继续指定文字、配色和风格。"

    # 更高分辨率
    fig = plt.figure(figsize=(10, 17.78), facecolor=colors["top"], dpi=200)
    ax = plt.axes([0, 0, 1, 1])
    ax.set_axis_off()
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    # 高质量纵向渐变背景
    top = np.array(_hex_to_rgb(colors["top"]))
    bottom = np.array(_hex_to_rgb(colors["bottom"]))
    grad = np.linspace(0, 1, 1600)[:, None]
    rgb = top * (1 - grad) + bottom * grad
    bg = np.repeat(rgb[:, None, :], 1200, axis=1)
    ax.imshow(bg, extent=[0, 1, 0, 1], aspect="auto", interpolation="bilinear")

    # 装饰性图形元素
    _draw_decorative_shapes(ax, style, colors)

    # 获取布局
    layout = _get_layout_template(style, colors)

    # 顶部装饰线
    ax.plot([0.25, 0.75], [layout["line_y"], layout["line_y"]],
            transform=ax.transAxes, color=colors["accent"], lw=2.5, alpha=0.85)

    # 顶部风格标签
    tag = style or "AI POSTER"
    ax.text(
        0.5, layout["tag_y"],
        tag,
        transform=ax.transAxes,
        ha="center", va="center",
        fontsize=16, color=colors["muted"],
        alpha=0.92, fontweight="bold",
        fontfamily="sans-serif",
    )

    # 主标题 - 更好的换行处理
    title_width = 10
    title_lines = textwrap.wrap(title, width=title_width)
    title_text = "\n".join(title_lines)
    title_fontsize = layout["title_size"]
    if len(title) > 12:
        title_fontsize = 36
    elif len(title) > 8:
        title_fontsize = 40

    ax.text(
        0.5, layout["title_y"],
        title_text,
        transform=ax.transAxes,
        ha="center", va="center",
        fontsize=title_fontsize,
        color=colors["text"],
        fontweight="bold",
        linespacing=1.3,
        fontfamily="sans-serif",
    )

    # 副标题
    subtitle_lines = textwrap.wrap(subtitle, width=18)
    subtitle_text = "\n".join(subtitle_lines)
    ax.text(
        0.5, layout["subtitle_y"],
        subtitle_text,
        transform=ax.transAxes,
        ha="center", va="center",
        fontsize=layout["subtitle_size"],
        color=colors["muted"],
        linespacing=1.4,
        fontfamily="sans-serif",
    )

    # 分隔线
    ax.plot([0.35, 0.65], [layout["subtitle_y"] - 0.04, layout["subtitle_y"] - 0.04],
            transform=ax.transAxes, color=colors["accent"], lw=1.5, alpha=0.5)

    # 视觉元素标签
    if elements:
        start_y = layout["elements_start_y"]
        for idx, element in enumerate(elements):
            y = start_y - idx * layout["element_spacing"]
            # 背景卡片
            ax.add_patch(plt.FancyBboxPatch(
                (0.20, y - 0.02), 0.60, 0.035,
                boxstyle="round,pad=0.005",
                transform=ax.transAxes,
                color="#ffffff", alpha=0.1, lw=0
            ))
            # 文字
            ax.text(
                0.5, y,
                f"✦ {element}",
                transform=ax.transAxes,
                ha="center", va="center",
                fontsize=17,
                color=colors["text"],
                alpha=0.95,
                fontfamily="sans-serif",
            )

    # 底部文案
    body_lines = textwrap.wrap(body, width=20)
    body_text = "\n".join(body_lines)
    ax.text(
        0.5, layout["body_y"],
        body_text,
        transform=ax.transAxes,
        ha="center", va="center",
        fontsize=layout["body_size"],
        color=colors["muted"],
        linespacing=1.6,
        fontfamily="sans-serif",
    )

    # 底部品牌标识
    ax.plot([0.40, 0.60], [layout["brand_y"] + 0.02, layout["brand_y"] + 0.02],
            transform=ax.transAxes, color=colors["text"], lw=1, alpha=0.3)
    ax.text(
        0.5, layout["brand_y"],
        "AI WORK AGENT",
        transform=ax.transAxes,
        ha="center", va="center",
        fontsize=13,
        color=colors["text"],
        alpha=0.5,
        fontweight="bold",
        letter_spacing=2,
        fontfamily="sans-serif",
    )

    # 高质量输出
    buffer = io.BytesIO()
    plt.savefig(buffer, format="png", dpi=200, bbox_inches="tight", pad_inches=0.02,
                facecolor=colors["top"])
    plt.close()
    buffer.seek(0)
    return buffer.getvalue()


def _decode_data_url(data_url: str) -> bytes:
    if "," not in data_url:
        raise ValueError("图片 data URL 格式不正确")
    return base64.b64decode(data_url.split(",", 1)[1])


def _extract_image_payload(result: dict[str, Any]) -> tuple[str, str]:
    """从常见图片生成响应中提取图片来源.

    Returns:
        (kind, value)，kind 为 url、data_url 或 b64。
    """
    data = result.get("data")
    if isinstance(data, list) and data:
        first = data[0] if isinstance(data[0], dict) else {}
        if first.get("b64_json"):
            return "b64", str(first["b64_json"])
        if first.get("url"):
            url = str(first["url"])
            return ("data_url", url) if url.startswith("data:image/") else ("url", url)

    output = result.get("output")
    if isinstance(output, dict):
        results = output.get("results")
        if isinstance(results, list) and results:
            first = results[0] if isinstance(results[0], dict) else {}
            if first.get("url"):
                return "url", str(first["url"])
            if first.get("b64_json"):
                return "b64", str(first["b64_json"])

    if result.get("url"):
        return "url", str(result["url"])
    if result.get("b64_json"):
        return "b64", str(result["b64_json"])

    raise ValueError("API 返回中未找到图片 URL 或 base64 数据")


async def generate_image(
    prompt: str,
    title: str = "AI 作图",
    subtitle: str = "",
    body: str = "",
    style: str = "",
    palette: str = "blue",
    elements: list[str] | None = None,
) -> bytes:
    """
    调用图片生成 API 生成图片

    Args:
        prompt: 图片描述提示词

    Returns:
        图片的二进制数据
    """
    image_url = settings.image_api_url
    image_key = settings.image_api_key or settings.anthropic_api_key
    image_model = settings.image_model

    if not image_url and not image_key:
        return generate_prompt_image(prompt, title, subtitle, body, style, palette, elements)

    if not image_url:
        # 根据 anthropic_base_url 自动推断
        base = settings.anthropic_base_url.rstrip("/")
        if not base:
            return generate_prompt_image(prompt, title, subtitle, body, style, palette, elements)
        if "/anthropic" in base:
            base = base.split("/anthropic")[0]
        image_url = f"{base}/v1/images/generations"

    try:
        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(
                image_url,
                headers={
                    "Authorization": f"Bearer {image_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": image_model,
                    "prompt": prompt,
                    "n": 1,
                    "size": "1024x1024",
                },
            )

        if resp.status_code == 200:
            result = resp.json()
            payload_kind, payload_value = _extract_image_payload(result)

            if payload_kind == "b64":
                return base64.b64decode(payload_value)
            if payload_kind == "data_url":
                return _decode_data_url(payload_value)

            # 下载图片
            async with httpx.AsyncClient(timeout=60) as client:
                img_resp = await client.get(payload_value)
                if img_resp.status_code == 200:
                    return img_resp.content
                else:
                    raise ValueError(f"下载图片失败: {img_resp.status_code}")
        else:
            logger.error("图片生成 API 错误: %s %s", resp.status_code, resp.text)
            raise ValueError(f"图片生成失败 ({resp.status_code}): {resp.text}")

    except Exception as e:
        logger.exception("调用图片生成 API 失败")
        raise
